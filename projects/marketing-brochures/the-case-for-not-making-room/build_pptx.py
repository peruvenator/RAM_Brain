"""
build_pptx.py
Generates The-Case-for-NOT-Making-Room-for-Alternatives.pptx as a
pixel-close replication of the 7-page HTML brochure.

Canvas: 8.5 x 11 in (US Letter portrait).
Units: everything measured in points (1 in = 72 pt).

Run: python build_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────────────────────────────────────

ROOT       = Path(__file__).parent.resolve()
BRAND      = ROOT / ".." / ".." / ".." / "references" / "brand-assets" / "return-stacked"
LOGO_BLACK = BRAND / "Logos" / "RS_Portfolio_Solutions_logos" / "Return Stacked Portfolio Solutions Black.png"
LOGO_WHITE = BRAND / "Logos" / "RS_Portfolio_Solutions_logos" / "Return Stacked Portfolio Solutions White.png"
BACKDROP   = BRAND / "Background_images" / "RS-Background-Blue.png"
CONTACT_BG = ROOT / "extracted_images" / "image3.jpg"

CHART_GROWTH    = ROOT / "chart_exports" / "chart_growth.png"
CHART_PIES      = ROOT / "chart_exports" / "chart_pies.png"
CHART_STACKEDBAR = ROOT / "chart_exports" / "chart_stackedbar.png"

OUT = ROOT / "The-Case-for-NOT-Making-Room-for-Alternatives.pptx"

for p in [LOGO_BLACK, LOGO_WHITE, BACKDROP, CONTACT_BG, CHART_GROWTH, CHART_PIES, CHART_STACKEDBAR]:
    if not p.exists():
        raise FileNotFoundError(f"Asset missing: {p}")

# ─────────────────────────────────────────────────────────────────────────────
# Brand tokens (from HTML CSS custom properties)
# ─────────────────────────────────────────────────────────────────────────────

C = {
    "text_primary":   RGBColor(0x2c, 0x36, 0x41),
    "text_secondary": RGBColor(0x62, 0x5c, 0x6d),
    "cover_dark":     RGBColor(0x17, 0x2c, 0x3a),
    "section_gray":   RGBColor(0xf0, 0xf1, 0xf1),
    "white":          RGBColor(0xff, 0xff, 0xff),
    "black":          RGBColor(0x00, 0x00, 0x00),
    "teal_primary":   RGBColor(0x14, 0xcf, 0xa6),
    "teal_light":     RGBColor(0xa1, 0xd7, 0xc6),
    "blue_secondary": RGBColor(0x3a, 0x6a, 0x9c),
    "blue_light":     RGBColor(0x7d, 0xa5, 0xce),
    "border_gray":    RGBColor(0xbf, 0xbf, 0xbf),
    "yellow":         RGBColor(0xeb, 0xe9, 0x6a),
    "teal_tint":      RGBColor(0xe8, 0xfa, 0xf4),  # rgba(20,207,166,0.06) on white ~ approx
}

FONT = "DM Sans"

# Page dimensions in points
PAGE_W = 8.5 * 72  # 612
PAGE_H = 11  * 72  # 792

# ─────────────────────────────────────────────────────────────────────────────
# Presentation scaffolding
# ─────────────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(8.5)
prs.slide_height = Inches(11)
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def pt(x): return Pt(x)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    """Add a rectangle shape at (x,y) with size (w,h). All values in pt."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pt(x), pt(y), pt(w), pt(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = pt(line_w)
    shp.text_frame.margin_left   = pt(0)
    shp.text_frame.margin_right  = pt(0)
    shp.text_frame.margin_top    = pt(0)
    shp.text_frame.margin_bottom = pt(0)
    return shp

def add_picture(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), pt(x), pt(y), pt(w), pt(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), pt(x), pt(y), width=pt(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), pt(x), pt(y), height=pt(h))
    return slide.shapes.add_picture(str(path), pt(x), pt(y))

def add_text(slide, x, y, w, h, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             margin=(0, 0, 0, 0), line_spacing=None, fill=None):
    """Add a textbox. `runs` is a list of dicts: {text, size, bold, italic, color, font}.
    margin is (top, right, bottom, left) in pt."""
    shp = slide.shapes.add_textbox(pt(x), pt(y), pt(w), pt(h))
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_top, tf.margin_right, tf.margin_bottom, tf.margin_left = \
        pt(margin[0]), pt(margin[1]), pt(margin[2]), pt(margin[3])
    tf.vertical_anchor = anchor

    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()

    # Support multiple paragraphs via runs containing a `newline` marker
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing

    for i, r in enumerate(runs):
        if r.get("newline"):
            p = tf.add_paragraph()
            p.alignment = r.get("align", align)
            if line_spacing is not None:
                p.line_spacing = line_spacing
            if r.get("space_before") is not None:
                p.space_before = pt(r["space_before"])
            if r.get("space_after") is not None:
                p.space_after = pt(r["space_after"])
            continue
        run = p.add_run()
        run.text = r["text"]
        f = run.font
        f.name  = r.get("font", FONT)
        f.size  = pt(r.get("size", 10))
        f.bold  = r.get("bold", False)
        f.italic = r.get("italic", False)
        col = r.get("color", C["text_primary"])
        f.color.rgb = col
    return shp

def simple_text(slide, x, y, w, h, text, size=10, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=None, margin=(0,0,0,0), fill=None):
    color = color or C["text_primary"]
    return add_text(slide, x, y, w, h,
                    [{"text": text, "size": size, "bold": bold, "italic": italic, "color": color}],
                    align=align, anchor=anchor, margin=margin,
                    line_spacing=line_spacing, fill=fill)

def teal_accent_bar(slide, x, y, h):
    """The 6pt teal bar that sits to the left of section headings."""
    return add_rect(slide, x, y, 6, h, fill=C["teal_primary"])

# ─────────────────────────────────────────────────────────────────────────────
# Header / Footer (shared across interior pages)
# ─────────────────────────────────────────────────────────────────────────────

def draw_interior_header(slide):
    """Header ~74pt tall: logo left, title block right, 10pt teal accent bar on right edge."""
    H = 74  # header height in pt
    # Logo
    # Logo image is ~38pt tall
    logo_h = 38
    logo_x = 36
    logo_y = (H - logo_h) / 2
    add_picture(slide, LOGO_BLACK, logo_x, logo_y, h=logo_h)

    # Title block (right-aligned)
    # Series title: 11pt teal bold
    # Topic title:  11pt text-primary medium
    # The whole block sits in the header region; right padding 18pt + 10pt accent bar
    accent_w = 10
    pad_right = 18
    title_block_w = 260
    title_block_x = PAGE_W - accent_w - pad_right - title_block_w
    title_block_y = (H - 32) / 2
    # Series
    simple_text(slide, title_block_x, title_block_y, title_block_w, 14,
                "The Case for NOT Making Room",
                size=11, bold=True, color=C["teal_primary"],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
    # Topic
    simple_text(slide, title_block_x, title_block_y + 14, title_block_w, 14,
                "for Alternatives",
                size=11, bold=False, color=C["text_primary"],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

    # Right-edge teal accent bar (full header height)
    add_rect(slide, PAGE_W - accent_w, 0, accent_w, H, fill=C["teal_primary"])


def draw_interior_footer(slide):
    """Footer: cover-dark icon square | teal band | white url area."""
    FOOTER_H = 28
    y = PAGE_H - FOOTER_H
    # Icon square
    add_rect(slide, 0, y, FOOTER_H, FOOTER_H, fill=C["cover_dark"])
    # Band (teal) fills middle
    band_x = FOOTER_H
    url_w = 180  # white url area width
    band_w = PAGE_W - FOOTER_H - url_w
    add_rect(slide, band_x, y, band_w, FOOTER_H, fill=C["teal_primary"])
    # URL area (white)
    add_rect(slide, band_x + band_w, y, url_w, FOOTER_H, fill=C["white"])
    simple_text(slide, band_x + band_w, y, url_w, FOOTER_H,
                "www.returnstacked.com",
                size=9, bold=True, color=C["text_primary"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ─────────────────────────────────────────────────────────────────────────────
# Section heading (teal accent bar + bold text), returns heading bottom y
# ─────────────────────────────────────────────────────────────────────────────

def draw_section_heading(slide, x, y, w, text, size=13, bar_h=14):
    """Draws a section heading at (x,y) with teal left-bar and bold text."""
    # Text sits at y, bar is centered vertically on text line (approx 14pt tall)
    # Bar at x, text starts at x + 14
    pad_left = 14
    line_h = size * 1.25
    # Bar vertical-center aligned with text
    bar_y = y + (line_h - bar_h) / 2
    add_rect(slide, x, bar_y, 6, bar_h, fill=C["teal_primary"])
    simple_text(slide, x + pad_left, y, w - pad_left, line_h + 4,
                text, size=size, bold=True, color=C["text_primary"],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2)
    return y + line_h


# ─────────────────────────────────────────────────────────────────────────────
# Dark panel (teal heading bar + gray body)
# ─────────────────────────────────────────────────────────────────────────────

def draw_dark_panel(slide, x, y, w, heading, bullets, heading_bg=None):
    """Teal-headed panel with gray body of bullets. Returns bottom y."""
    heading_bg = heading_bg or C["teal_primary"]
    # Heading bar
    head_h = 22  # 11pt text + ~9pt top + 9pt bottom padding (actually 22pt total)
    add_rect(slide, x, y, w, head_h, fill=heading_bg)
    simple_text(slide, x + 16, y, w - 32, head_h,
                heading, size=11, bold=True, color=C["white"],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Body
    body_pad_top = 12
    body_pad_bot = 12
    body_pad_x   = 18
    bullet_size  = 10
    line_sp      = 1.55
    line_h = bullet_size * line_sp
    # Rough body height: pad + (bullets * line_h with wrap allowance)
    # We approximate; caller may want exact sizing. We'll generate a rect then textbox.
    # Compute approximate height based on text length per bullet
    import math
    # Rough chars-per-line at ~10pt with w-36 px width
    approx_cpl = max(30, int((w - 2*body_pad_x) / 4.2))
    total_lines = 0
    for b in bullets:
        total_lines += max(1, math.ceil(len(b) / approx_cpl))
    body_h = body_pad_top + body_pad_bot + total_lines * line_h + (len(bullets) - 1) * 5
    add_rect(slide, x, y + head_h, w, body_h, fill=C["section_gray"])
    # Text frame with bullets
    runs = []
    for i, b in enumerate(bullets):
        if i > 0:
            runs.append({"newline": True, "space_before": 2})
        runs.append({"text": "•  " + b, "size": bullet_size, "color": C["text_primary"]})
    add_text(slide, x, y + head_h, w, body_h, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             margin=(body_pad_top, body_pad_x, body_pad_bot, body_pad_x),
             line_spacing=line_sp)
    return y + head_h + body_h


# ─────────────────────────────────────────────────────────────────────────────
# Callout teal (italic teal quote w/ left border)
# ─────────────────────────────────────────────────────────────────────────────

def draw_callout_teal(slide, x, y, w, text, size=12):
    line_sp = 1.4
    import math
    approx_cpl = max(30, int((w - 20) / 4.0))
    lines = max(1, math.ceil(len(text) / approx_cpl))
    h = lines * size * line_sp + 12
    # Left border
    add_rect(slide, x, y, 3, h, fill=C["teal_primary"])
    simple_text(slide, x + 16, y, w - 16, h, text,
                size=size, bold=True, italic=True, color=C["teal_primary"],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=line_sp)
    return y + h


def draw_callout_teal_wide(slide, x, y, w, text, size=14):
    """Full-width callout with padding and subtle tinted background."""
    line_sp = 1.4
    import math
    approx_cpl = max(30, int((w - 44) / 4.7))
    lines = max(1, math.ceil(len(text) / approx_cpl))
    h = lines * size * line_sp + 36
    # Background tint
    add_rect(slide, x, y, w, h, fill=C["teal_tint"])
    # Left border (4pt teal)
    add_rect(slide, x, y, 4, h, fill=C["teal_primary"])
    simple_text(slide, x + 24, y, w - 48, h, text,
                size=size, bold=True, italic=True, color=C["teal_primary"],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=line_sp)
    return y + h


# ─────────────────────────────────────────────────────────────────────────────
# Gray section box
# ─────────────────────────────────────────────────────────────────────────────

def draw_gray_section(slide, x, y, w, runs, pad=(16, 18, 16, 18), line_spacing=1.55):
    # Estimate height from runs' total text length
    import math
    total_chars = sum(len(r.get("text", "")) for r in runs if not r.get("newline"))
    paragraph_count = sum(1 for r in runs if r.get("newline")) + 1
    approx_cpl = max(30, int((w - pad[1] - pad[3]) / 4.2))
    lines_base = max(paragraph_count, math.ceil(total_chars / approx_cpl))
    body_h = pad[0] + pad[2] + lines_base * 10 * line_spacing + (paragraph_count - 1) * 4
    add_rect(slide, x, y, w, body_h, fill=C["section_gray"])
    add_text(slide, x, y, w, body_h, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             margin=pad, line_spacing=line_spacing)
    return y + body_h


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════

def build_cover():
    s = prs.slides.add_slide(BLANK)

    # Cover top bar (white, 96pt)
    top_h = 96
    add_rect(s, 0, 0, PAGE_W, top_h, fill=C["white"])
    # Logo (48pt tall) inside, 40pt from left
    logo_h = 48
    add_picture(s, LOGO_BLACK, 40, (top_h - logo_h) / 2, h=logo_h)

    # Cover banner: backdrop image + teal strip + title
    banner_y = top_h
    banner_h = 340
    # Backdrop
    add_picture(s, BACKDROP, 0, banner_y, w=PAGE_W, h=banner_h)
    # Dark navy overlay approximation — skipped; backdrop already colored.
    # 10pt teal accent bar on left
    add_rect(s, 0, banner_y, 10, banner_h, fill=C["teal_primary"])

    # Title block (left-aligned)
    title_x = 54
    title_y = banner_y + 110
    # "The Case for NOT" — 36pt semibold, "NOT" bold 700
    # Line 1
    add_text(s, title_x, title_y, 520, 50,
             [{"text": "The Case for ", "size": 36, "bold": True, "color": C["white"]},
              {"text": "NOT",            "size": 36, "bold": True, "color": C["white"]}],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.1)
    # Line 2: "Making Room for"
    simple_text(s, title_x, title_y + 48, 520, 50, "Making Room for",
                size=36, bold=True, color=C["white"],
                align=PP_ALIGN.LEFT, line_spacing=1.1)
    # Line 3: "Alternatives" — 40pt bold teal
    simple_text(s, title_x, title_y + 98, 520, 60, "Alternatives",
                size=40, bold=True, color=C["teal_primary"],
                align=PP_ALIGN.LEFT, line_spacing=1.1)

    # Cover intro: heading, text, highlight
    intro_y = banner_y + banner_h + 40
    simple_text(s, 56, intro_y, PAGE_W - 112, 30,
                "Rethinking How Diversification Gets Done",
                size=20, bold=True, color=C["text_primary"],
                align=PP_ALIGN.CENTER, line_spacing=1.25)

    intro_text_y = intro_y + 40
    simple_text(s, 56, intro_text_y, PAGE_W - 112, 90,
                "Most portfolios don't struggle with theory. They struggle with implementation. "
                "Advisors know diversification matters, but adding alternatives often means selling "
                "core positions, introducing unfamiliar line items, and defending positions that "
                "lag at the wrong time.",
                size=12, bold=False, color=C["text_primary"],
                align=PP_ALIGN.CENTER, line_spacing=1.6)

    hl_y = intro_text_y + 92
    simple_text(s, 56, hl_y, PAGE_W - 112, 30,
                "What if diversification didn't require giving something up?",
                size=18, bold=True, color=C["teal_primary"],
                align=PP_ALIGN.CENTER, line_spacing=1.25)

    # Cover meta (bottom)
    meta_y = PAGE_H - 22 - 10
    simple_text(s, 40, meta_y, 300, 12,
                "© Return Stacked® Portfolio Solutions, 2026",
                size=8, color=C["text_secondary"],
                align=PP_ALIGN.LEFT)
    simple_text(s, PAGE_W - 40 - 260, meta_y, 260, 12,
                "For Investment Professional Use Only",
                size=8, color=C["text_secondary"],
                align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 2 — THEORY vs REALITY + THE REAL TRADE-OFF
# ═════════════════════════════════════════════════════════════════════════════

def build_page2():
    s = prs.slides.add_slide(BLANK)
    draw_interior_header(s)

    content_x = 36
    content_w = PAGE_W - 72

    # Section 1 heading
    y = 88
    y = draw_section_heading(s, content_x, y, content_w,
                             "Theory vs. Reality: Why Diversification Is Harder Than It Should Be")
    # Body
    y += 12
    body1 = ("Most advisors understand the value of diversification. The academic case is clear: "
             "combining assets with different return drivers can improve risk-adjusted returns over time. "
             "But knowing what to own is only half the battle. The real challenge is making diversification "
             "work inside actual portfolios, real models, and real client conversations.")
    simple_text(s, content_x, y, content_w, 60, body1,
                size=10, color=C["text_primary"], line_spacing=1.55)
    y += 64

    # Growth chart
    chart_h = 180
    add_picture(s, CHART_GROWTH, content_x, y, w=content_w, h=chart_h)
    y += chart_h + 8

    # Source note
    src1 = ("Source: Bloomberg; Credit Suisse; Société Générale. U.S. Stocks is the S&P 500 Total "
            "Return Index (SPX). U.S. Bonds is the Bloomberg US Aggregate Bond Index (LBUSTRUU). "
            "CTA Trend is the Société Générale Trend Index (NEIXCTAT). Merger Arbitrage is the LAB "
            "Merger Arbitrage Liquid Index (CSLABMA). Gold is spot gold quoted in US Dollars "
            "(XAU Currency). You cannot invest in an index. Returns are gross of fees and taxes. "
            "Past performance is not indicative of future results. Period is 12/31/1999 through 12/31/2024.")
    simple_text(s, content_x, y, content_w, 32, src1,
                size=7, italic=True, color=C["text_secondary"], line_spacing=1.3)
    y += 34

    # Section 2 heading
    y = draw_section_heading(s, content_x, y, content_w,
                             "The Real Trade-Off: The Uncomfortable Reality")
    y += 10
    body2 = ("To add diversifiers, most portfolios must reduce core stock and bond exposure, accept "
             "tracking error versus benchmarks, introduce new line items clients don't recognize, "
             "and defend positions that behave differently, especially at the wrong time.")
    simple_text(s, content_x, y, content_w, 48, body2,
                size=10, color=C["text_primary"], line_spacing=1.55)
    y += 52

    # Pie charts
    pie_h = 140
    add_picture(s, CHART_PIES, content_x, y, w=content_w, h=pie_h)
    y += pie_h + 8

    # Source note
    src2 = ('"60/40" refers to a 60% allocation to stocks and a 40% allocation to bonds. '
            '"50/30/20" refers to a 50% allocation to stocks, a 30% allocation to bonds, and '
            'a 20% allocation to alternatives.')
    simple_text(s, content_x, y, content_w, 18, src2,
                size=7, italic=True, color=C["text_secondary"], line_spacing=1.3)
    y += 22

    # Callout teal wide — full content width
    draw_callout_teal_wide(s, content_x, y, content_w,
        "Good diversification ideas often fail because of how they're implemented, not why they exist.")

    draw_interior_footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 3 — THE REAL CONSTRAINT
# ═════════════════════════════════════════════════════════════════════════════

def build_page3():
    s = prs.slides.add_slide(BLANK)
    draw_interior_header(s)

    content_x = 36
    content_w = PAGE_W - 72

    y = 88
    y = draw_section_heading(s, content_x, y, content_w,
        "The Real Constraint: What If the Problem Isn't the Diversifier?")
    y += 12
    body1 = ("Diversifiers frequently get blamed for underperformance in portfolios. But the issue "
             "may not be the strategy itself. It may be how it's implemented. The behavioral and "
             "structural pressures of holding standalone alternative positions can undermine even "
             "the best-designed allocation.")
    simple_text(s, content_x, y, content_w, 60, body1,
                size=10, color=C["text_primary"], line_spacing=1.55)
    y += 70

    # Dark panel
    y = draw_dark_panel(s, content_x, y, content_w,
        "The Behavioral Challenge of Standalone Diversifiers",
        [
            "They can lag during strong equity markets, creating opportunity cost anxiety",
            "They experience uncomfortable drawdowns when viewed in isolation",
            'They become the focus of "why do we own this?" conversations',
            "They get reduced or removed, often right before they're needed most",
        ])
    y += 14

    body2 = ("These are not flaws in the diversifiers. They are structural consequences of how "
             "diversifiers are typically held: as visible, standalone line items that compete for "
             "attention with core positions. Sometimes the issue isn't with what you own, it's how "
             "you own it.")
    simple_text(s, content_x, y, content_w, 72, body2,
                size=10, color=C["text_primary"], line_spacing=1.55)
    y += 80

    # Gray section
    gray_runs = [
        {"text": "The pattern is familiar: ", "size": 10, "bold": True, "color": C["text_primary"]},
        {"text": ("an advisor adds a well-reasoned diversifier to client portfolios. It works exactly "
                  "as intended, zigging when stocks zag. But during the next bull market run, the "
                  "diversifier lags. Clients question it. The advisor, under pressure, trims or "
                  "removes it. Then the next downturn arrives, and the protection is gone. The "
                  "diversifier didn't fail. The implementation did."),
         "size": 10, "color": C["text_primary"]},
    ]
    draw_gray_section(s, content_x, y, content_w, gray_runs)

    draw_interior_footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 4 — INTRODUCING RETURN STACKING
# ═════════════════════════════════════════════════════════════════════════════

def build_page4():
    s = prs.slides.add_slide(BLANK)
    draw_interior_header(s)

    content_x = 36
    content_w = PAGE_W - 72

    y = 88
    y = draw_section_heading(s, content_x, y, content_w,
        "Introducing Return Stacking: A Structural Approach to Diversification")
    y += 12
    body1 = ("Some investors approach diversification as a portfolio-structure problem: combining "
             "exposures so diversification doesn't require giving something up. This framework, "
             "often called return stacking, seeks to improve diversification without increasing "
             "behavioral pressures by focusing on how exposures are combined at the portfolio level.")
    simple_text(s, content_x, y, content_w, 60, body1,
                size=10, color=C["text_primary"], line_spacing=1.55)
    y += 70

    # Stacked bar chart (tall, max-width 50%)
    chart_h = 260
    chart_w = content_w * 0.5
    chart_x = content_x + (content_w - chart_w) / 2
    add_picture(s, CHART_STACKEDBAR, chart_x, y, w=chart_w, h=chart_h)
    y += chart_h + 8

    # Source note
    src = ("For illustrative purposes only. The red line represents the cutoff between a traditional "
           "portfolio exposure and the exposure an investor may receive through a return stacked portfolio.")
    simple_text(s, content_x, y, content_w, 20, src,
                size=7, italic=True, color=C["text_secondary"], line_spacing=1.3)
    y += 24

    # Section 2 heading
    y = draw_section_heading(s, content_x, y, content_w,
        "Structure changes the experience, not the objective")
    y += 14

    # Two-column comparison
    col_gap = 20
    col_w = (content_w - col_gap) / 2
    # Left column: Traditional (teal head)
    y_left = draw_dark_panel(s, content_x, y, col_w, "Traditional Approach", [
        "Diversifiers appear as standalone line items and draw scrutiny",
        "Allocation decisions become an either-or choice between beta and diversifiers",
        "Creates pressure to defend positions when they lag",
    ])
    # Right column: Return Stacked (cover-dark head)
    y_right = draw_dark_panel(s, content_x + col_w + col_gap, y, col_w, "Return Stacked Approach", [
        "Diversifiers are paired with core beta in a single line item",
        'Allocation decisions become "yes, and": holding diversifiers without sacrificing beta',
        "Reduces pressure to explain standalone performance",
    ], heading_bg=C["cover_dark"])

    draw_interior_footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 5 — OFFENSE + IMPLEMENTATION
# ═════════════════════════════════════════════════════════════════════════════

def build_page5():
    s = prs.slides.add_slide(BLANK)
    draw_interior_header(s)

    content_x = 36
    content_w = PAGE_W - 72

    y = 88
    y = draw_section_heading(s, content_x, y, content_w,
        "Beyond Defense: Return Stacking for Offense")
    y += 12
    # Body with mixed formatting: "trying to be better than the benchmark, can we simply add to it?"
    add_text(s, content_x, y, content_w, 40,
        [
            {"text": "Most investors think of diversifiers as defensive tools. Return stacking "
                     "reframes the question: instead of trying to be ", "size": 10},
            {"text": "better", "size": 10, "italic": True},
            {"text": " than the benchmark, can we simply ", "size": 10},
            {"text": "add to it", "size": 10, "bold": True},
            {"text": "?", "size": 10},
        ],
        align=PP_ALIGN.LEFT, line_spacing=1.55)
    y += 44

    simple_text(s, content_x, y, content_w, 18,
                "There are three fundamental approaches to pursuing outperformance:",
                size=10, line_spacing=1.55)
    y += 22

    # Approach cards: 3 columns
    card_gap = 14
    card_w = (content_w - 2 * card_gap) / 3
    card_h = 120
    headers = [
        ("Security Selection",  C["blue_secondary"],
         "Can we pick better securities than the benchmark?",
         "Relies on stock-picking skill to identify undervalued positions or avoid overvalued ones."),
        ("Tactical Allocation", C["blue_light"],
         "Can we time exposures better than the benchmark?",
         "Relies on timing skill to shift between asset classes at the right moment."),
        ("Return Stacking",     C["teal_primary"],
         "Can we add return streams on top of the benchmark?",
         "Relies on structural efficiency to layer additional exposures without sacrificing core allocations."),
    ]
    for i, (title, bg, q, a) in enumerate(headers):
        cx = content_x + i * (card_w + card_gap)
        cy = y
        # Card container (border)
        add_rect(s, cx, cy, card_w, card_h, fill=C["white"], line=C["border_gray"], line_w=1)
        # Header bar
        head_h = 24
        add_rect(s, cx, cy, card_w, head_h, fill=bg)
        simple_text(s, cx + 14, cy, card_w - 28, head_h, title,
                    size=10, bold=True, color=C["white"],
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        body_runs = [
            {"text": q, "size": 9, "italic": True, "color": C["text_secondary"]},
            {"newline": True, "space_before": 8},
            {"text": a, "size": 9, "color": C["text_primary"]},
        ]
        add_text(s, cx, cy + head_h, card_w, card_h - head_h, body_runs,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 margin=(14, 14, 14, 14), line_spacing=1.5)
    y += card_h + 24

    # Second section heading
    y = draw_section_heading(s, content_x, y, content_w,
        "How We Think About Implementation")
    y += 12
    simple_text(s, content_x, y, content_w, 18,
                "If this structural approach resonates, we typically start the conversation by exploring three questions:",
                size=10, color=C["text_primary"], line_spacing=1.55)
    y += 22

    # Gray section with 3 bullets
    bullets = [
        ("Where has diversification been hardest to hold?",
         " Understanding past pain points reveals where structural solutions may help most."),
        ("Which trade-offs matter most in your models?",
         " Not every portfolio has the same constraints. Identifying the binding ones helps focus the conversation."),
        ("Would a structural approach help, or not?",
         " Return stacking isn't the answer to every portfolio question. Honest evaluation matters."),
    ]
    # Build runs
    runs = []
    for i, (q, a) in enumerate(bullets):
        if i > 0:
            runs.append({"newline": True, "space_before": 6})
        runs.append({"text": "•  ", "size": 10, "color": C["text_primary"]})
        runs.append({"text": q, "size": 10, "bold": True, "color": C["text_primary"]})
        runs.append({"text": a, "size": 10, "color": C["text_primary"]})
    y2 = draw_gray_section(s, content_x, y, content_w, runs)
    y = y2 + 14

    # Closing callout
    draw_callout_teal(s, content_x, y, content_w,
                      "It's not about adding more. It's about adding smarter.")

    draw_interior_footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 6 — CONTACT
# ═════════════════════════════════════════════════════════════════════════════

def build_page6():
    s = prs.slides.add_slide(BLANK)
    # Full-bleed backdrop
    add_picture(s, CONTACT_BG, 0, 0, w=PAGE_W, h=PAGE_H)
    # Dark overlay (~45% black)
    overlay = add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=C["black"])
    overlay.fill.fore_color.rgb = C["black"]
    # Set transparency via XML (python-pptx lacks direct API)
    sppr = overlay.fill._xPr
    solidfill = sppr.find(qn("a:solidFill"))
    if solidfill is not None:
        srgb = solidfill.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", "45000")  # 45%

    # White logo top-left
    add_picture(s, LOGO_WHITE, 44, 36, h=42)

    # Title "Learn More" big white
    simple_text(s, 60, 320, PAGE_W - 120, 80,
                "Learn More",
                size=52, bold=True, color=C["white"],
                align=PP_ALIGN.LEFT, line_spacing=1.1)
    # URL teal
    simple_text(s, 60, 410, PAGE_W - 120, 30,
                "www.returnstacked.com",
                size=20, bold=True, color=C["teal_primary"],
                align=PP_ALIGN.LEFT)
    # Contact Us button (teal rect with white bold text)
    btn_x, btn_y, btn_w, btn_h = 60, 460, 150, 36
    add_rect(s, btn_x, btn_y, btn_w, btn_h, fill=C["teal_primary"])
    simple_text(s, btn_x, btn_y, btn_w, btn_h, "Contact Us",
                size=13, bold=True, color=C["white"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom meta
    simple_text(s, 44, PAGE_H - 32, 300, 16,
                "© Return Stacked® Portfolio Solutions, 2026",
                size=8, color=C["white"],
                align=PP_ALIGN.LEFT)
    simple_text(s, PAGE_W - 44 - 260, PAGE_H - 32, 260, 16,
                "For Investment Professional Use Only",
                size=8, color=C["white"],
                align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# PAGE 7 — DISCLOSURES
# ═════════════════════════════════════════════════════════════════════════════

def build_page7():
    s = prs.slides.add_slide(BLANK)
    # Black background
    add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=C["black"])

    content_x = 36
    content_w = PAGE_W - 72
    y = 40
    # Top header bar
    simple_text(s, content_x, y, 300, 16,
                "Return Stacked® Portfolio Solutions",
                size=10, bold=True, color=C["white"],
                align=PP_ALIGN.LEFT)
    simple_text(s, PAGE_W - content_x - 100, y, 100, 16, "2026",
                size=10, color=RGBColor(0x80, 0x80, 0x80),
                align=PP_ALIGN.RIGHT)
    y += 20
    # Border-bottom approximation
    add_rect(s, content_x, y, content_w, 0.75,
             fill=RGBColor(0x33, 0x33, 0x33))
    y += 14

    # Heading
    simple_text(s, content_x, y, content_w, 22,
                "General Disclosures",
                size=16, bold=True, color=C["white"],
                align=PP_ALIGN.LEFT)
    y += 26

    paragraphs = [
        ('The information set forth in this document has been obtained or derived from sources believed '
         'by Newfound Research LLC ("Newfound") to be reliable. However, Newfound does not make any '
         'representation or warranty, express or implied, as to the information\'s accuracy or '
         'completeness, nor does Newfound recommend that the information serve as the basis of any '
         'investment decision.'),
        ('Certain information contained in this document constitutes "forward-looking statements," '
         'which can be identified by the use of forward-looking terminology such as "may," "will," '
         '"should," "expect," "anticipate," "project," "estimate," "intend," "continue," or "believe," '
         'or the negatives thereof or other variations or comparable terminology. Due to various risks '
         'and uncertainties, actual events or results or the actual performance of an investment '
         'managed using any of the investment strategies or styles described in this document may '
         'differ materially from those reflected in such forward-looking statements. The information '
         'in this document is made available on an "as is," without representation or warranty basis.'),
        ('There can be no assurance that any investment strategy or style will achieve any level of '
         'performance, and investment results may vary substantially from year to year or even from '
         'month to month. An investor could lose all or substantially all of his or her investment. '
         'Both the use of a single adviser and the focus on a single investment strategy could result '
         'in the lack of diversification and consequently, higher risk. The information herein is not '
         'intended to provide, and should not be relied upon for, accounting, legal or tax advice or '
         'investment recommendations. Any investment strategy and themes discussed herein may be '
         'unsuitable for investors depending on their specific investment objectives and financial '
         'situation. You should consult your investment adviser, tax, legal, accounting or other '
         'advisors about the matters discussed herein. These materials represent an assessment of the '
         'market environment at specific points in time and are intended neither to be a guarantee of '
         'future events nor as a primary basis for investment decisions. Past performance is not '
         'indicative of future performance and investments in equity securities do present risk of loss.'),
        ('Investors should understand that while performance results may show a general rising trend '
         'at times, there is no assurance that any such trends will continue. If such trends are broken, '
         'then investors may experience real losses. The information included in this presentation '
         'reflects the different assumptions, views and analytical methods of Newfound as of the date '
         'of this document. The views expressed reflect the current views as of the date hereof and '
         'neither the author nor Newfound undertakes to advise you of any changes in the views expressed herein.'),
        ('This presentation has been provided solely for informational purposes and does not constitute '
         'a current or past recommendation or an offer or solicitation of an offer, or any advice or '
         'recommendation, to purchase any securities or other financial instruments, and may not be '
         'construed as such. This presentation should not be considered as investment advice or a '
         'recommendation of any particular security, strategy or investment product.'),
        ('Return stacking may involve the use of derivatives, leverage, and short selling, each of '
         'which may increase potential losses and risk.'),
        ('No part of this document may be reproduced in any form, or referred to in any other '
         'publication, without express written permission from Newfound Research.'),
    ]

    # Build a single large text frame with all paragraphs
    runs = []
    for i, para in enumerate(paragraphs):
        if i > 0:
            runs.append({"newline": True, "space_before": 7})
        runs.append({"text": para, "size": 7.5,
                     "color": RGBColor(0xcc, 0xcc, 0xcc)})
    add_text(s, content_x, y, content_w, PAGE_H - y - 60, runs,
             align=PP_ALIGN.LEFT, line_spacing=1.4)

    # Copyright + case number at bottom
    simple_text(s, content_x, PAGE_H - 50, content_w, 14,
                "© Return Stacked® Portfolio Solutions, 2026. All rights reserved.",
                size=7.5, color=RGBColor(0x99, 0x99, 0x99),
                align=PP_ALIGN.LEFT)
    simple_text(s, content_x, PAGE_H - 34, content_w, 14,
                "Case #1c89b077-b462-4766-8a31-4d046eee8774",
                size=7.5, color=RGBColor(0x66, 0x66, 0x66),
                align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    build_cover()
    build_page2()
    build_page3()
    build_page4()
    build_page5()
    build_page6()
    build_page7()
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
