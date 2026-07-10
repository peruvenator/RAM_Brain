"""
Lazy Money Flyer -- PPTX Build Script
Generates a 2-slide US Letter (8.5 x 11) PowerPoint flyer
matching the HTML version with ReSolve AM brand assets.

Usage:
    python build_flyer_pptx.py
"""

import base64
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
REPO = HERE.parent.parent
BRAND = REPO / "references" / "brand-assets" / "resolve-am"
FONTS_DIR = BRAND / "Fonts" / "Helvetica Neue LT Std"
CHARTS = HERE / "chart_exports"
OUTPUT = HERE / "lazy-money-flyer.pptx"

# ---------------------------------------------------------------------------
# Brand colors
# ---------------------------------------------------------------------------
NAVY       = RGBColor(0x03, 0x2F, 0x69)
PRIMARY    = RGBColor(0x00, 0x47, 0x8D)
MED_BLUE   = RGBColor(0x47, 0xA3, 0xDA)
SKY        = RGBColor(0x89, 0xD2, 0xFF)
ICE        = RGBColor(0x8F, 0xCE, 0xF1)
AMBER      = RGBColor(0xFB, 0xBA, 0x00)
TEXT_BLUE   = RGBColor(0x29, 0x4A, 0x85)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
BORDER     = RGBColor(0xB5, 0xC4, 0xCE)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Font name (falls back if not installed)
FONT = "Helvetica Neue LT Std"
FONT_FALLBACK = "Helvetica Neue"


def decode_b64_to_stream(b64_path: Path) -> BytesIO:
    """Decode a base64 text file to a BytesIO image stream."""
    raw = b64_path.read_text().strip()
    buf = BytesIO(base64.b64decode(raw))
    buf.seek(0)
    return buf


def png_stream(path: Path) -> BytesIO:
    """Read a PNG file into a BytesIO stream."""
    buf = BytesIO(path.read_bytes())
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height, text, *,
                font_size=10, bold=False, italic=False,
                color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                font_name=FONT, line_spacing=1.15):
    """Add a simple single-run text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, *,
                          font_size=10, bold=False, color=DARK_GRAY,
                          alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, *,
                    font_size=9.5, color=DARK_GRAY):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        p.line_spacing = Pt(font_size * 1.5)
        p.level = 0
        # Bullet character
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        # Set bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        pPr.append(buChar)
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': '00478D'})
        buClr.append(srgb)
        pPr.append(buClr)
        buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})
        pPr.append(buSzPct)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, *, line_color=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_banner(slide, top=0):
    """Add the geometric blue banner strip across the top."""
    W = 8.5
    H = 0.42
    # Base navy bar
    add_rect(slide, 0, top, W, H, NAVY)
    # Geometric triangle overlays using freeform shapes
    t = top
    b = top + H

    # Triangle 1: medium blue, left portion
    specs = slide.shapes.build_freeform(Inches(0), Inches(t))
    specs.add_line_segments([
        (Inches(2.2), Inches(t)),
        (Inches(1.4), Inches(b)),
        (Inches(0), Inches(b)),
    ])
    tri1 = specs.convert_to_shape()
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = MED_BLUE
    tri1.line.fill.background()

    # Triangle 2: sky blue, left-center
    specs = slide.shapes.build_freeform(Inches(1.0), Inches(t))
    specs.add_line_segments([
        (Inches(2.8), Inches(t)),
        (Inches(2.0), Inches(b)),
        (Inches(0.2), Inches(b)),
    ])
    tri2 = specs.convert_to_shape()
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = SKY
    tri2.line.fill.background()

    # Triangle 3: ice blue, center
    specs = slide.shapes.build_freeform(Inches(2.4), Inches(t))
    specs.add_line_segments([
        (Inches(4.0), Inches(t)),
        (Inches(3.2), Inches(b)),
        (Inches(1.6), Inches(b)),
    ])
    tri3 = specs.convert_to_shape()
    tri3.fill.solid()
    tri3.fill.fore_color.rgb = ICE
    tri3.line.fill.background()

    # Triangle 4: medium blue, right side
    specs = slide.shapes.build_freeform(Inches(5.5), Inches(t))
    specs.add_line_segments([
        (Inches(7.0), Inches(t)),
        (Inches(6.2), Inches(b)),
        (Inches(4.7), Inches(b)),
    ])
    tri4 = specs.convert_to_shape()
    tri4.fill.solid()
    tri4.fill.fore_color.rgb = MED_BLUE
    tri4.line.fill.background()

    # Triangle 5: navy dark accent, far right
    specs = slide.shapes.build_freeform(Inches(6.5), Inches(t))
    specs.add_line_segments([
        (Inches(8.5), Inches(t)),
        (Inches(8.5), Inches(b)),
        (Inches(5.7), Inches(b)),
    ])
    tri5 = specs.convert_to_shape()
    tri5.fill.solid()
    tri5.fill.fore_color.rgb = PRIMARY
    tri5.line.fill.background()


def add_table(slide, left, top, width, height):
    """Add the stats comparison table."""
    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = table_shape.table

    # Column widths
    tbl.columns[0].width = Inches(width * 0.30)
    tbl.columns[1].width = Inches(width * 0.22)
    tbl.columns[2].width = Inches(width * 0.24)
    tbl.columns[3].width = Inches(width * 0.24)

    # Header row
    headers = ["Statistics", "Balanced\nPortfolio", "Carry\n(Excess Returns)", "100% Balanced\n+ 100% Carry"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        run.font.name = FONT
        run.font.size = Pt(7.5)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    data = [
        ["Annualized Return",     "4.92%",   "11.34%",  "17.90%"],
        ["Annualized Volatility", "11.14%",  "22.08%",  "20.43%"],
        ["Max Drawdown",          "-21.23%", "-23.98%", "-14.71%"],
        ["Sharpe Ratio",          "0.44",    "0.51",    "0.88"],
    ]
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.name = FONT
            run.font.size = Pt(8.5)
            run.font.color.rgb = PRIMARY if j == 3 else DARK_GRAY
            run.font.bold = (j == 3)
            if j == 0:
                run.font.color.rgb = TEXT_BLUE
                run.font.bold = False
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Remove default table borders and style
    from pptx.oxml.ns import qn
    tbl_xml = tbl._tbl
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('bandRow', '0')
        tblPr.set('firstRow', '0')

    return table_shape


# ---------------------------------------------------------------------------
# Slide 1 -- Strategy & Benefits
# ---------------------------------------------------------------------------
def build_slide1(slide):
    # Banner strip
    add_banner(slide, top=0)

    # Logo
    logo_stream = decode_b64_to_stream(BRAND / "cover-template" / "logo-black_b64.txt")
    slide.shapes.add_picture(logo_stream, Inches(0.5), Inches(0.65), Inches(2.0))

    # Headline
    add_textbox(slide, 0.5, 1.25, 7.5, 1.0,
                "Put Your Lazy Money\nto Work",
                font_size=36, bold=True, color=NAVY,
                line_spacing=1.1)

    # Subtitle
    add_textbox(slide, 0.5, 2.15, 7.5, 0.3,
                "PORTABLE ALPHA FOR THE MID-MARKET",
                font_size=12, bold=False, color=PRIMARY)

    # Divider line
    add_rect(slide, 0.5, 2.52, 7.5, 0.015, BORDER)

    # Portfolio Enhancement Diagram (slide 4)
    slide4_img = png_stream(CHARTS / "slide4.png")
    slide.shapes.add_picture(slide4_img, Inches(0.5), Inches(2.65), Inches(7.5))

    # Two-column section below diagram
    col_top = 5.85

    # Left column: What is Return Stacking?
    add_textbox(slide, 0.5, col_top, 3.5, 0.25,
                "WHAT IS RETURN STACKING?",
                font_size=10, bold=True, color=PRIMARY)

    add_textbox(slide, 0.5, col_top + 0.3, 3.5, 1.5,
                "Return stacking is a capital-efficient investment approach that "
                "allows investors to maintain their existing portfolio while adding "
                "uncorrelated return streams on top. By using futures-based strategies "
                "funded with T-Bill collateral, investors can achieve more than 100% "
                "economic exposure without additional leverage on the underlying portfolio.",
                font_size=8.5, color=DARK_GRAY, line_spacing=1.45)

    # Right column: Why ReSolve?
    add_textbox(slide, 4.5, col_top, 3.5, 0.25,
                "WHY RESOLVE?",
                font_size=10, bold=True, color=PRIMARY)

    add_bullet_list(slide, 4.5, col_top + 0.3, 3.5, 2.0, [
        "20+ years of quantitative research heritage",
        "Systematic carry, trend, and volatility strategies across 50+ global futures markets",
        "True portable alpha: uncorrelated to stocks and bonds",
        "Institutional infrastructure, mid-market accessibility",
        "Separately managed accounts with full transparency",
    ], font_size=8.5, color=DARK_GRAY)

    # Pull quote
    quote_top = 7.65

    # Amber accent bar
    add_rect(slide, 0.5, quote_top, 0.04, 0.65, AMBER)

    # Quote background
    add_rect(slide, 0.54, quote_top, 7.46, 0.65, LIGHT_GRAY)

    # Quote text
    add_textbox(slide, 0.75, quote_top + 0.06, 7.0, 0.35,
                "\u201CReSolve has delivered institutional-grade alpha that has "
                "genuinely enhanced our portfolio construction.\u201D",
                font_size=9.5, italic=True, color=TEXT_BLUE,
                line_spacing=1.4)

    add_textbox(slide, 0.75, quote_top + 0.42, 7.0, 0.2,
                "\u2014 Jonathan Glidden, Former CIO, Delta Air Lines",
                font_size=8, bold=False, color=DARK_GRAY)

    # Corner decoration (subtle, low opacity)
    corner_stream = decode_b64_to_stream(BRAND / "cover-template" / "corner-decoration_b64.txt")
    pic = slide.shapes.add_picture(corner_stream, Inches(5.2), Inches(8.6), Inches(3.3))
    # Set transparency via XML
    from pptx.oxml.ns import qn
    blipFill = pic._element.find(qn('p:blipFill'))
    if blipFill is not None:
        blip = blipFill.find(qn('a:blip'))
        if blip is not None:
            alphaModFix = blip.makeelement(qn('a:alphaModFix'), {'amt': '8000'})
            blip.append(alphaModFix)

    # Footer URL
    add_textbox(slide, 0.5, 10.5, 7.5, 0.3,
                "www.investresolve.com",
                font_size=10, bold=False, color=PRIMARY,
                alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 -- Proof & CTA
# ---------------------------------------------------------------------------
def build_slide2(slide):
    # Banner strip
    add_banner(slide, top=0)

    # Dark navy header section
    add_rect(slide, 0, 0.42, 8.5, 0.85, NAVY)

    add_textbox(slide, 0.5, 0.48, 7.5, 0.5,
                "Building Blocks for a\nRobust Stacked Portfolio",
                font_size=22, bold=True, color=WHITE,
                line_spacing=1.1)

    add_textbox(slide, 0.5, 1.0, 7.5, 0.25,
                "Case Study on a Live ReSolve Carry Mandate",
                font_size=10, bold=False, color=SKY)

    # Annual Returns Chart
    annual_img = png_stream(CHARTS / "annual_returns.png")
    slide.shapes.add_picture(annual_img, Inches(0.5), Inches(1.45), Inches(7.5))

    # Correlation Heatmap (left) + Stats Table (right)
    corr_img = png_stream(CHARTS / "correlation_heatmap.png")
    slide.shapes.add_picture(corr_img, Inches(0.5), Inches(4.6), Inches(3.4))

    # Stats table
    add_table(slide, 4.1, 4.6, 4.0, 1.8)

    # Institutional Credibility badges
    cred_top = 7.0
    cred_items = [
        "$2B+ in AUM across strategies",
        "Sub-advised for major ETF platforms",
        "CFTC-registered CTA",
        "SEC-registered RIA",
    ]
    cred_width = 1.8
    cred_gap = 0.15
    for i, item in enumerate(cred_items):
        x = 0.5 + i * (cred_width + cred_gap)
        # Background
        add_rect(slide, x, cred_top, cred_width, 0.4, LIGHT_GRAY)
        # Left accent bar
        add_rect(slide, x, cred_top, 0.04, 0.4, PRIMARY)
        # Text
        add_textbox(slide, x + 0.12, cred_top + 0.05, cred_width - 0.15, 0.3,
                    item, font_size=8, bold=False, color=TEXT_BLUE)

    # CTA box
    cta_top = 7.65
    add_rect(slide, 0.5, cta_top, 7.5, 0.7, NAVY)

    add_textbox(slide, 0.7, cta_top + 0.08, 5.0, 0.3,
                "Ready to put your lazy money to work?",
                font_size=13, bold=True, color=WHITE)

    add_textbox(slide, 0.7, cta_top + 0.38, 5.0, 0.2,
                "info@investresolve.com  |  www.investresolve.com",
                font_size=9, bold=False, color=SKY)

    # Logo in CTA (white version)
    logo_white_stream = decode_b64_to_stream(BRAND / "cover-template" / "logo-white_b64.txt")
    slide.shapes.add_picture(logo_white_stream, Inches(6.2), cta_top + Inches(0.12), Inches(1.6))

    # Source note
    add_textbox(slide, 0.5, 8.55, 7.5, 0.6,
                "Source: Tiingo. Analysis by ReSolve Asset Management SEZC (Cayman). "
                "ReSolve Carry is ReSolve Futures Yield (Carry) 20% Volatility Program (Excess Returns). "
                "Bonds is the iShares Core U.S. Aggregate Bond ETF (AGG). "
                "Global Equities is the iShares MSCI ACWI ETF (ACWI). "
                "Period is from September 1, 2021 through March 31, 2026. "
                "Indicated returns of one year or more are annualized. "
                "These results are a carveout of returns for the Program. "
                "All performance data is provided by the third party fund admin less a 0.85% annual fee.",
                font_size=5.5, color=RGBColor(0x88, 0x88, 0x88),
                line_spacing=1.35)

    # Disclosures
    add_textbox(slide, 0.5, 9.35, 7.5, 0.6,
                "PAST PERFORMANCE IS NOT A GUARANTEE OF FUTURE RESULTS. THE RISK OF LOSS IN "
                "TRADING COMMODITY INTERESTS IS SUBSTANTIAL. These materials do not constitute "
                "an offer or solicitation of an offer to make an investment in any of the funds "
                "or separately managed accounts ReSolve Global manages. ReSolve Global operates "
                "within a fund of funds and as excess returns (calculated prior to any yield on "
                "posted collateral). Confidential \u2014 Qualified Eligible Purchasers Only.",
                font_size=5, color=RGBColor(0x99, 0x99, 0x99),
                line_spacing=1.35)

    # Corner decoration
    corner_stream = decode_b64_to_stream(BRAND / "cover-template" / "corner-decoration_b64.txt")
    pic = slide.shapes.add_picture(corner_stream, Inches(5.2), Inches(8.6), Inches(3.3))
    from pptx.oxml.ns import qn
    blipFill = pic._element.find(qn('p:blipFill'))
    if blipFill is not None:
        blip = blipFill.find(qn('a:blip'))
        if blip is not None:
            alphaModFix = blip.makeelement(qn('a:alphaModFix'), {'amt': '8000'})
            blip.append(alphaModFix)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build():
    print("Building PPTX...")

    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]

    # Slide 1
    slide1 = prs.slides.add_slide(blank_layout)
    build_slide1(slide1)
    print("  Slide 1 (Strategy & Benefits) done")

    # Slide 2
    slide2 = prs.slides.add_slide(blank_layout)
    build_slide2(slide2)
    print("  Slide 2 (Proof & CTA) done")

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
